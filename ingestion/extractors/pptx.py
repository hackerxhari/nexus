"""
PowerPoint extractor — handles .pptx and .ppt files.
Extracts slide text, tables, notes, and shapes preserving structure.
Uses python-pptx (already in requirements).
"""

import os
from typing import List

from pptx import Presentation
from pptx.util import Inches

from core.exceptions import ExtractionFailedError
from core.logger import get_logger
from ingestion.extractors.base import BaseExtractor, ExtractionResult

logger = get_logger(__name__)


class PPTXExtractor(BaseExtractor):

    @property
    def supported_extensions(self) -> List[str]:
        return ["pptx", "ppt"]

    def _extract(self, filepath: str) -> ExtractionResult:
        if not os.path.exists(filepath):
            raise ExtractionFailedError(
                os.path.basename(filepath),
                "File not found"
            )

        warnings = []
        slide_sections = []

        try:
            prs = Presentation(filepath)
            slide_count = len(prs.slides)

            if slide_count == 0:
                raise ExtractionFailedError(
                    os.path.basename(filepath),
                    "Presentation has no slides"
                )

            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text_parts = []
                slide_text_parts.append(f"\n[Slide {slide_num}]")

                # Extract text from all shapes
                for shape in slide.shapes:
                    try:
                        # Text frames (titles, body text, text boxes)
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                text = paragraph.text.strip()
                                if text:
                                    slide_text_parts.append(text)

                        # Tables — extract with row-column relationships
                        if shape.has_table:
                            table_text = self._extract_table(shape.table)
                            if table_text:
                                slide_text_parts.append(
                                    f"[Table]\n{table_text}"
                                )

                    except Exception as e:
                        warnings.append(
                            f"Slide {slide_num}, shape error: {str(e)}"
                        )
                        continue

                # Extract slide notes
                try:
                    if slide.has_notes_slide:
                        notes_text = slide.notes_slide.notes_text_frame.text
                        if notes_text and notes_text.strip():
                            slide_text_parts.append(
                                f"[Notes] {notes_text.strip()}"
                            )
                except Exception as e:
                    warnings.append(
                        f"Slide {slide_num}, notes error: {str(e)}"
                    )

                # Only add slide if it has content beyond the header
                if len(slide_text_parts) > 1:
                    slide_sections.append("\n".join(slide_text_parts))

        except ExtractionFailedError:
            raise
        except Exception as e:
            raise ExtractionFailedError(
                os.path.basename(filepath),
                f"PowerPoint parsing failed: {str(e)}"
            )

        combined_text = "\n\n".join(slide_sections)

        if not combined_text.strip():
            warnings.append("No text content found in presentation")

        return ExtractionResult(
            text=combined_text,
            page_count=slide_count,
            extraction_method="pptx_parser",
            warnings=warnings,
            metadata={
                "slide_count": slide_count,
                "slides_with_content": len(slide_sections)
            }
        )

    def _extract_table(self, table) -> str:
        """
        Convert a PowerPoint table to readable text.
        Uses first row as headers for row-column relationship structuring.
        """
        rows = []
        headers = []

        for row_idx, row in enumerate(table.rows):
            cells = [cell.text.strip() for cell in row.cells]

            if row_idx == 0:
                # First row as headers
                headers = cells
                # Only output header row if it has content
                if any(c for c in headers):
                    rows.append(" | ".join(c for c in headers if c))
                continue

            # For data rows: format as Header: Value pairs
            if headers and any(c for c in cells):
                pairs = []
                for col_idx, cell_text in enumerate(cells):
                    if cell_text:
                        header = headers[col_idx] if col_idx < len(headers) else f"Col{col_idx + 1}"
                        if header:
                            pairs.append(f"{header}: {cell_text}")
                        else:
                            pairs.append(cell_text)
                if pairs:
                    rows.append(" | ".join(pairs))
            elif any(c for c in cells):
                # No headers — just join cells
                row_text = " | ".join(c for c in cells if c)
                if row_text:
                    rows.append(row_text)

        return "\n".join(rows)
